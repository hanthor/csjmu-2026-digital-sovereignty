#!/usr/bin/env python3
# /// script
# requires-python = ">=3.11"
# dependencies = ["python-pptx", "Pillow"]
# ///
"""Build a PPTX from a Typst file by rasterising each page as a PNG slide."""

import argparse
import subprocess
import sys
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Emu


def build_pptx(source: str, output: str, ppi: int = 200) -> None:
    with tempfile.TemporaryDirectory() as tmp:
        png_pattern = str(Path(tmp) / "slide-{p}.png")
        result = subprocess.run(
            ["typst", "compile", source, png_pattern, "--ppi", str(ppi)],
            capture_output=True, text=True,
        )
        if result.returncode != 0:
            print(result.stderr, file=sys.stderr)
            sys.exit(1)

        pages = sorted(
            Path(tmp).glob("slide-*.png"),
            key=lambda p: int(p.stem.split("-")[1]),
        )
        if not pages:
            print("No pages produced.", file=sys.stderr)
            sys.exit(1)

        # Determine slide dimensions from first page
        with Image.open(pages[0]) as im:
            px_w, px_h = im.size

        # Convert pixels → EMU (English Metric Units used by PPTX)
        emu_w = Emu(int(px_w / ppi * 914400))
        emu_h = Emu(int(px_h / ppi * 914400))

        prs = Presentation()
        prs.slide_width  = emu_w
        prs.slide_height = emu_h

        blank_layout = prs.slide_layouts[6]  # completely blank

        for png in pages:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(str(png), 0, 0, emu_w, emu_h)

        prs.save(output)
        print(f"Written {len(pages)} slides → {output}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("source", help="Input .typ file")
    parser.add_argument("output", help="Output .pptx file")
    parser.add_argument("--ppi", type=int, default=200, help="Render resolution (default 200)")
    args = parser.parse_args()
    build_pptx(args.source, args.output, args.ppi)
